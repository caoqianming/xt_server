from pptx import Presentation
from pptx.util import Inches
from server.settings import BASE_DIR
from apps.audit.models import Atask, AtaskIssue, AtaskTeam
import os
from datetime import datetime
from pptx.util import Pt
from PIL import Image
from apps.system.models import User
from apps.utils.permission import has_perm
from rest_framework.exceptions import ParseError
from .models import R_LEVEL_DICT
from io import BytesIO
import logging

myLogger = logging.getLogger("log")

templ = os.path.join(BASE_DIR, "media/muban/安全审计总结.pptx")

def convert_to_supported_format(image_path):
    """将不支持的图片格式（如 MPO）转换为 PNG/JPEG"""
    try:
        with Image.open(image_path) as img:
            # 创建一个临时内存文件（避免磁盘 I/O）
            output_buffer = BytesIO()
            
            # 如果图片是 MPO，Pillow 会自动读取第一帧
            if img.format == "MPO":
                myLogger.warning(f"Converted MPO image to PNG: {image_path}")
            
            # 转换为 PNG（或 JPEG，根据需求调整）
            img.save(output_buffer, format="PNG")
            output_buffer.seek(0)  # 重置指针
            return output_buffer
    except Exception as e:
        myLogger.error(f"Failed to convert image {image_path}: {str(e)}")
        raise ValueError(f"Unsupported image format or corrupt file: {image_path}")

def add_image_to_slide(slide, image_path, left, top, width=None, height=None):
    """安全地添加图片到幻灯片，自动处理格式转换"""
    try:
        # 直接尝试插入图片（如果格式已支持）
        slide.shapes.add_picture(image_path, left, top, width, height)
    except ValueError as e:
        if "unsupported image format" in str(e):
            # 如果是格式不支持（如 MPO），转换后再插入
            converted_image = convert_to_supported_format(image_path)
            slide.shapes.add_picture(converted_image, left, top, width, height)
        else:
            raise  # 其他错误继续抛出

def export_pptx(atask:Atask, FileName:str, user:User):
   
    def replace_text_preserve_formatting(shape, old_text, new_text):
        if not shape.has_text_frame:
            return False
        
        text_frame = shape.text_frame
        found = False
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    # 替换文本但保留run的格式
                    run.text = run.text.replace(old_text, new_text)
                    found = True
        
        return found
    ex_type = "user"
    if has_perm(user, ["atask.update"]):
        ex_type = "all"
    elif AtaskTeam.objects.filter(atask=atask, member=user).exists():
        ex_type = "user"
    else:
        raise ParseError("没有导出PPT的权限")
    prs = Presentation(templ)

    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.text == "company_name":
                replace_text_preserve_formatting(shape, "company_name", atask.company.name)
            elif shape.text == "date":
                replace_text_preserve_formatting(shape, "date", datetime.now().strftime("%Y-%m-%d"))
            elif shape.text == "leader":
                replace_text_preserve_formatting(shape, "leader", user.name)
    
    if ex_type == "all":
        issues = AtaskIssue.objects.filter(atask=atask).order_by("atask", "standarditem__number_sort", "-create_time")
    elif ex_type == "user":
        issues = AtaskIssue.objects.filter(atask=atask, create_by=user).order_by("atask", "standarditem__number_sort", "-create_time")
    
    issues = issues.select_related('standarditem', 'photos').prefetch_related('photos')
    for ind, issue in enumerate(issues):
        photos = issue.photos

        left = top = Inches(1.2)
        if photos.exists():
            top = Inches(1.2)
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        
        # 3. 将新幻灯片移动到指定位置
        target_position = ind + 5
        slides = prs.slides._sldIdLst  # 获取幻灯片列表
        slides.insert(target_position, slides[-1])  # 将最后一张插入到目标位置

        shapes = slide.shapes

        cols = 3
        rows = 2
        width = Inches(11.4)
        height = Inches(1.2)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(7.4)

        # write column headings
        table.cell(0, 0).text = '条款'
        table.cell(0, 1).text = '问题风险等级'
        table.cell(0, 2).text = '问题描述'

        # write body cells
        standarditem = issue.standarditem
        if standarditem:
            if standarditem.level == 30:
                yj = standarditem.parent.parent.content
            elif standarditem.level == 20:
                yj = standarditem.parent.content
            elif standarditem.level == 10:
                yj = standarditem.content
        table.cell(1, 0).text = f'{standarditem.number} {yj}'  if standarditem else ""
        table.cell(1, 1).text = R_LEVEL_DICT.get(issue.risk_level, "")
        table.cell(1, 2).text = issue.content

        # 设置不同的标题和内容字体大小
        for row in range(rows):
            for col in range(cols):
                tf = table.cell(row, col).text_frame
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        if row == 0:  # 标题行
                            run.font.size = Pt(18)  # 标题字体大小
                            run.font.bold = True    # 标题加粗
                        else:         # 内容行
                            # if col == 2:  # 第3列（索引从0开始，所以2表示第3列）
                            #     run.font.size = Pt(18)  # 第3列字体缩小
                            # else:
                                run.font.size = Pt(20)  # 其他内容列字体大小

        if photos.exists():
            num_images = min(photos.count(), 3)  # 不超过3张
            img_height = Inches(4.2)  # 固定高度
            min_width = Inches(2.4) 
            spacing = Inches(0.2)     # 图片间隔
            top_position = Inches(2.8)  # 图片起始Y位置
        
        if issue.content and len(issue.content) > 48:
            top_position = Inches(3.2)  # 图片起始Y位置
            img_height = Inches(3.8)  # 固定高度

            # 计算所有图片的自适应宽度
            img_widths = []
            for v in photos.all()[:3]:
                img_path = BASE_DIR + v.path
                with Image.open(img_path) as img:
                    width_px, height_px = img.size
                    aspect_ratio = width_px / height_px  # 宽高比（宽度/高度）h)
                    # width = max(img_height * aspect_ratio, min_width)
                width = img_height * aspect_ratio
                img_widths.append(width)

            # 计算总宽度和起始位置（居中）
            total_width = sum(img_widths) + (num_images - 1) * spacing
            start_left = (Inches(10) - total_width) / 2 + Inches(2)  # 居中

            # 添加图片
            current_left = start_left
            for i, v in enumerate(photos.all()[:3]):
                img_path = BASE_DIR + v.path
                add_image_to_slide(slide, img_path, current_left, top_position, img_widths[i], img_height)
                # slide.shapes.add_picture(
                #     img_path,
                #     current_left, top_position,
                #     height=img_height,
                #     width=img_widths[i]
                # )
                current_left += img_widths[i] + spacing  # 移动到下一张图片位置

    path = f"/media/temp/{FileName}_{atask.id}_{user.name}.pptx"
    prs.save(BASE_DIR+path)
    return path